# -*- coding: utf-8 -*-
"""QueryStrategist 功能介绍 PPT 生成脚本（Step 0–2 检索策略包版）

依赖 python-pptx；若未安装，用 _shared_tools/scripts/ensure_tool.py 安装。
输出到脚本所在目录（相对路径，不硬编码私人路径）。
"""
import os
import argparse
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ---------- palette ----------
PRIMARY   = RGBColor(0x1F, 0x4E, 0x79)
TEAL      = RGBColor(0x2E, 0x7D, 0x8B)
ACCENT    = RGBColor(0xE8, 0x83, 0x3A)
LIGHT     = RGBColor(0xF2, 0xF5, 0xF9)
DARK      = RGBColor(0x22, 0x2A, 0x33)
GRAY      = RGBColor(0x6B, 0x76, 0x83)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
SOFTBLUE  = RGBColor(0xCF, 0xDD, 0xEC)
PALEBLUE  = RGBColor(0xE3, 0xE9, 0xF0)
CREAM     = RGBColor(0xE8, 0xD8, 0xC8)
FONT = "微软雅黑"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

def set_cjk(run, name=FONT):
    run.font.name = name
    rPr = run._r.get_or_add_rPr()
    latin = rPr.find(qn('a:latin'))
    if latin is None:
        latin = rPr.makeelement(qn('a:latin'), {})
        latin.set('typeface', name)
        rPr.insert(0, latin)
    latin.set('typeface', name)
    ea = rPr.find(qn('a:ea'))
    if ea is None:
        ea = rPr.makeelement(qn('a:ea'), {})
        ea.set('typeface', name)
        latin.addnext(ea)
    else:
        ea.set('typeface', name)

def style(run, size=18, bold=False, color=DARK, name=FONT):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    set_cjk(run, name)

def add_slide():
    return prs.slides.add_slide(BLANK)

def fill_shape(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.shadow.inherit = False

def bg(slide, color):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    fill_shape(sh, color)
    return sh

def add_footer(s, text="QueryStrategist · SCP 广场参赛作品（Step 0–2）"):
    tb = s.shapes.add_textbox(Inches(0.5), Inches(7.02), Inches(12.3), Inches(0.35))
    tf = tb.text_frame
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.RIGHT
    r = p.add_run(); r.text = text
    style(r, size=10, color=GRAY)

def content_slide(title, kicker=None):
    s = add_slide()
    bg(s, LIGHT)
    band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.15))
    fill_shape(band, PRIMARY)
    acc = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.18), Inches(1.15))
    fill_shape(acc, ACCENT)
    tb = s.shapes.add_textbox(Inches(0.5), Inches(0.12), Inches(12.3), Inches(0.95))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = title
    style(p.runs[0], size=28, bold=True, color=WHITE)
    if kicker:
        p2 = tf.add_paragraph(); p2.text = kicker
        style(p2.runs[0], size=13, color=SOFTBLUE)
    add_footer(s)
    return s

def bullets(slide, items, left=0.7, top=1.5, width=11.9, height=5.2, size=17, gap=10):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame; tf.word_wrap = True
    first = True
    for it in items:
        if len(it) == 3:
            text, lvl, color = it
        else:
            text, lvl = it; color = DARK
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.level = lvl
        p.space_after = Pt(gap)
        p.line_spacing = 1.12
        bullet_char = "•  " if lvl == 0 else "–  "
        r = p.add_run(); r.text = bullet_char + text
        style(r, size=size if lvl == 0 else size - 2, bold=(lvl == 0 and color != DARK), color=color)
    return tb

def placeholder(slide, left, top, width, height, text):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 Inches(left), Inches(top), Inches(width), Inches(height))
    box.fill.solid(); box.fill.fore_color.rgb = PALEBLUE
    box.line.color.rgb = GRAY; box.line.width = Pt(1)
    box.shadow.inherit = False
    tf = box.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text
    style(r, size=13, color=GRAY)

# ============ SLIDE 1 — Title ============
s = add_slide()
bg(s, PRIMARY)
acc = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(6.62), prs.slide_width, Inches(0.12))
fill_shape(acc, ACCENT)
tb = s.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(11.7), Inches(1.9))
tf = tb.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "QueryStrategist"
style(r, size=54, bold=True, color=WHITE)
p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
r2 = p2.add_run(); r2.text = "Step 0–2 人机协作 · 智能文献检索策略生成器"
style(r2, size=26, bold=True, color=SOFTBLUE)
tb2 = s.shapes.add_textbox(Inches(0.8), Inches(4.05), Inches(11.7), Inches(0.8))
tf2 = tb2.text_frame; tf2.paragraphs[0].alignment = PP_ALIGN.CENTER
r3 = tf2.paragraphs[0].add_run()
r3.text = "一句话模糊科研意图 → 6 大数据库精准检索式 + 文献候选清单，即拿即用"
style(r3, size=18, color=CREAM)
tb3 = s.shapes.add_textbox(Inches(0.8), Inches(5.45), Inches(11.7), Inches(0.7))
tf3 = tb3.text_frame; tf3.paragraphs[0].alignment = PP_ALIGN.CENTER
r4 = tf3.paragraphs[0].add_run()
r4.text = "面向综述 · 研究论著 · 学位论文 · 开题报告 · 基金申请 · 调研报告   |   SCP 广场参赛作品   |   2026.08"
style(r4, size=13, color=RGBColor(0xB8, 0xC8, 0xD8))

# ============ SLIDE 2 — Pain & Value ============
s = content_slide("为什么需要 QueryStrategist", "把文献检索的「策略设计」，封装为可被复用、可追溯的标准化能力")
bullets(s, [
    ("检索起步难：从“我想研究 X 在 Y 中的应用”到可执行的高级检索式，跨度大、靠人工拼关键词易漏易滥", 0),
    ("跨库构建重复：WoS / Scopus / IEEE / Google Scholar / CNKI / 万方语法各异，逐库构建耗时数天", 0),
    ("策略缺依据：候选文献常凭经验筛，缺少跨库收割的量化覆盖与 OA 状态佐证", 0),
    ("价值主张", 0, ACCENT),
    ("将检索策略前期工作封装为 1 个主 Skill + 11 个子模块，可独立调用也可串联", 1),
    ("人类掌舵、AI 执行：AI 规模化执行（关键词收敛、检索式生成、跨库收割），人类在决策门拍板", 1),
    ("最终交付检索策略包四件套：范围卡 + 6 库检索式 + 候选清单 + 使用说明", 1),
])

# ============ SLIDE 3 — Overview ============
s = content_slide("方案总览：状态机驱动的 3 步（Step 0–2）流水线", "单入口触发 · 全局上下文继承 · 决策门强制人工确认")
bullets(s, [
    ("单入口触发：用户输入“开始文献检索”，状态机顺序串联 3 个子 Skill（Step 0–2）", 0),
    ("全局 Pipeline Context 跨步继承，避免重复劳动、保持上下文连续", 0),
    ("3 个决策门（G0–G2）强制人工确认，AI 不越权做价值判断", 0),
    ("双通道检索", 0, TEAL),
    ("Search A：手工检索式（6 大库：WoS / Scopus / IEEE / Google Scholar / CNKI / 万方，每库查全式 A + 查准式 B）", 1),
    ("Search B：API 自动收割（OpenAlex 收割 + Crossref 按 DOI 验证）", 1),
    ("按写作类型调策略权重：综述查全 / 论著查准 / 开题基金新颖性", 0, ACCENT),
])

# ============ SLIDE 4 — Pipeline diagram (Step 0–2) ============
s = content_slide("3 步流水线：从配置到检索策略包", "每一步结束设决策门，AI 呈现事实、人类拍板")
steps = [
    ("0", "预检配置", "Setup Wizard"),
    ("1", "范围界定", "Scope Definer"),
    ("2", "双通道检索", "Search Strategist V1"),
]
cols, rows = 3, 1
bw, bh = 3.60, 1.90
gx, gy = 0.35, 0.36
total_w = cols * bw + (cols - 1) * gx
start_x = (13.333 - total_w) / 2
start_y = 1.75
for i, (num, name, sub) in enumerate(steps):
    r = i // cols; c = i % cols
    x = start_x + c * (bw + gx)
    y = start_y + r * (bh + gy)
    color = PRIMARY if i < 2 else ACCENT
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                             Inches(x), Inches(y), Inches(bw), Inches(bh))
    box.fill.solid(); box.fill.fore_color.rgb = color
    box.line.color.rgb = WHITE; box.line.width = Pt(1)
    box.shadow.inherit = False
    tf = box.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.1); tf.margin_right = Inches(0.1)
    p1 = tf.paragraphs[0]; p1.alignment = PP_ALIGN.CENTER
    r1 = p1.add_run(); r1.text = "STEP " + num
    style(r1, size=12, bold=True, color=SOFTBLUE)
    p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run(); r2.text = name
    style(r2, size=17, bold=True, color=WHITE)
    p3 = tf.add_paragraph(); p3.alignment = PP_ALIGN.CENTER
    r3 = p3.add_run(); r3.text = sub
    style(r3, size=10, color=RGBColor(0xDD, 0xE7, 0xF0))
placeholder(s, 0.55, 4.55, 12.2, 1.7, "终点交付：检索策略包四件套（范围卡 + 6 库检索式 + 文献候选清单 + 使用说明）\n检索式按写作类型调权重（查全 A / 查准 B / 新颖性）\n候选清单含 OA 状态与 DOI 链接，人机闸门 G2 确认后交付")

# ============ SLIDE 5 — Core capabilities (scripted) ============
s = content_slide("核心能力：已实装可运行脚本", "超越纯提示词——关键机械化环节已封装为 Python，开箱即用")
bullets(s, [
    ("检索式生成  query_generator.py（零依赖）— 6 大库高级检索式一键产出（broad / precise / angle_tech / angle_app / review）", 0),
    ("文献收割  harvest.py — OpenAlex / Crossref，含配额守卫（预算 / 429 熔断 / 缓存 / dry-run）", 0),
    ("工具安装  ensure_tool.py — 开源工具检测/隔离安装（镜像直连，绕过 venv 回滚）", 0),
    ("策略包模板  search_strategy_pack_template.md — 范围卡 / 检索式合集 / 候选清单 / 使用说明标准结构", 0, ACCENT),
])

# ============ SLIDE 6 — Tech stack ============
s = content_slide("技术实现与工具链", "零外部模型依赖 · 直接引用不改写 · 全链路可追溯")
bullets(s, [
    ("收割：OpenAlex / Crossref 公开 API（网络；Crossref 仅按 DOI 做一致性验证）", 0),
    ("质量参考：SJR 期刊分级（CC BY-NC 4.0，用户自行下载，不随仓库分发）", 0),
    ("导出：openpyxl（MIT）文档与表格", 0),
    ("检索式：6 平台高级检索语法（WoS / Scopus / IEEE / Google Scholar / CNKI / 万方）", 0),
    ("设计原则：Step 0–2 零外部模型，仅用内置 LLM；检索式与候选清单全链路绑定上游出处", 0, ACCENT),
])
placeholder(s, 8.7, 1.6, 4.1, 4.6, "演示截图位\n（建议放置：检索策略包交付界面\n或某库检索式输出，如 scope_card）")

# ============ SLIDE 7 — Innovation ============
s = content_slide("创新与差异化", "为什么它不是又一个“检索工具”")
bullets(s, [
    ("铁律约束：“收割 ≠ 语料”——AI 呈现候选清单，下载决策权完全在用户，杜绝编造引用", 0, ACCENT),
    ("决策门强制人工确认（G0–G2），杜绝 AI 越权替人类决定纳入哪些文献", 0),
    ("按写作类型调检索策略权重（综述查全 / 论著查准 / 开题基金新颖性）——比数据库自带 Query Builder 强的核心差异", 0),
    ("可追溯证据链：检索策略包每个字段标注上游出处（范围卡 = 检索式依据 = 候选清单口径）", 0),
    ("零外部模型依赖（Step 0–2 仅内置 LLM），开箱即用、易部署", 0),
    ("1 个主 Skill + 11 个子模块完全自包含，含可运行 scripts/，可整体也可单步使用", 0),
])

# ============ SLIDE 8 — Scenarios ============
s = content_slide("应用场景与落地路径", "从科研到教学，覆盖六类文献写作")
bullets(s, [
    ("综述 / 系统综述：查全优先，宽式 A 为主，近 10 年时间窗", 0),
    ("研究论著 / 实验研究：查准优先，精准式 B 为主，近 5 年时间窗", 0, ACCENT),
    ("学位论文：查全+查准均衡，A+B 并重", 0),
    ("开题报告 / 基金申请：兼顾新颖性，近 2 年过滤 + 高被引标注", 0),
    ("调研报告：查全+查准均衡，A 起步可调窄", 0),
    ("落地方式：上架 SCP 广场，用户一句话触发完整 0–2 流水线，产出检索策略包", 0, TEAL),
])

# ============ SLIDE 9 — IP & compliance ============
s = content_slide("知识产权与合规", "开源可复用 · 第三方数据合规解耦")
bullets(s, [
    ("源码 MIT 开源，允许学习、修改、再分发（见 LICENSE）", 0),
    ("SJR 期刊质量数据集为 CC BY-NC 4.0（不可商用），不随仓库分发，用户自行下载", 0),
    ("1 个主 Skill + 11 个子模块完全自包含，目录遵循统一资源布局", 0),
    ("参赛亮点（面向专家评审）", 0, ACCENT),
    ("完整可运行：4 个脚本化能力 + 双通道检索 + 策略包四件套模板，非纯提示词", 1),
    ("方法学严谨：决策门 + 上游上下文继承 + 证据链，降低 AI 幻觉与越权", 1),
    ("六类写作类型适配 + 6 平台检索式生成：覆盖多种科研写作场景", 1),
])

# ============ SLIDE 10 — Closing ============
s = add_slide()
bg(s, PRIMARY)
acc = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(0.0), prs.slide_width, Inches(0.12))
fill_shape(acc, ACCENT)
tb = s.shapes.add_textbox(Inches(0.8), Inches(2.6), Inches(11.7), Inches(1.6))
tf = tb.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "即刻体验 QueryStrategist"
style(r, size=40, bold=True, color=WHITE)
p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
r2 = p2.add_run(); r2.text = "SCP 广场搜索 “QueryStrategist”，一句话开启你的检索策略生成"
style(r2, size=18, color=SOFTBLUE)
tb3 = s.shapes.add_textbox(Inches(0.8), Inches(5.6), Inches(11.7), Inches(0.6))
tf3 = tb3.text_frame; tf3.paragraphs[0].alignment = PP_ALIGN.CENTER
r3 = tf3.paragraphs[0].add_run()
r3.text = "感谢书生科学发现平台 · 期待你的收藏与反馈"
style(r3, size=14, color=CREAM)

# ---------- save ----------
# 默认输出到脚本所在目录；发布或评审时可显式指定输出路径。
parser = argparse.ArgumentParser(description="生成 QueryStrategist 功能介绍 PPT")
default_out = os.path.join(os.path.dirname(os.path.abspath(__file__)), "QueryStrategist_功能介绍.pptx")
parser.add_argument("--output", default=default_out, help="PPTX 输出路径")
args = parser.parse_args()
out_path = os.path.abspath(args.output)
os.makedirs(os.path.dirname(out_path), exist_ok=True)
prs.save(out_path)
print("saved:", out_path, "slides:", len(prs.slides._sldIdLst))
